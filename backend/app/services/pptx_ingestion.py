from __future__ import annotations

import hashlib
from dataclasses import dataclass
from io import BytesIO
from typing import Any
from uuid import UUID

import boto3
from sqlalchemy import delete, select

from app.core.settings import Settings, get_settings
from app.db.models.content_chunk import ContentChunk
from app.db.models.course_content import CourseContent
from app.db.models.document_page import DocumentPage
from app.db.session import get_session_maker


@dataclass(frozen=True)
class ExtractedSlide:
    slide_no: int
    title: str
    slide_text: str
    notes_text: str


def _s3_client(settings: Settings):
    kwargs: dict[str, Any] = {"service_name": "s3", "region_name": settings.s3_region}
    if settings.s3_endpoint_url:
        kwargs["endpoint_url"] = settings.s3_endpoint_url
    if settings.s3_access_key_id and settings.s3_secret_access_key:
        kwargs["aws_access_key_id"] = settings.s3_access_key_id
        kwargs["aws_secret_access_key"] = settings.s3_secret_access_key
    return boto3.client(**kwargs)


def _is_pptx(content: CourseContent) -> bool:
    mt = (content.mime_type or "").lower().strip()
    name = (content.original_filename or "").lower().strip()
    if "presentation" in mt or "powerpoint" in mt:
        return True
    return name.endswith(".pptx")


def _sha256_hex(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8", errors="ignore")).hexdigest()


def _text_splitter(settings: Settings):
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    chunk_size = int(getattr(settings, "rag_chunk_size", 1200))
    overlap = int(getattr(settings, "rag_chunk_overlap", 200))
    return RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=overlap)


def _join_nonempty(*parts: str, sep: str = "\n") -> str:
    items = []
    for p in parts:
        t = (p or "").strip()
        if t:
            items.append(t)
    return sep.join(items).strip()


def _extract_pptx_slides(pptx_bytes: bytes) -> list[ExtractedSlide]:
    """
    Extract slide text + speaker notes from a PPTX.

    NOTE: We import python-pptx lazily so the backend can still run in environments
    where it isn't installed; ingestion will no-op in that case.
    """
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return []

    prs = Presentation(BytesIO(pptx_bytes))

    out: list[ExtractedSlide] = []

    def _shape_text(shape) -> str:  # noqa: ANN001
        # Text frames and table cells.
        txt = ""
        try:
            if getattr(shape, "has_text_frame", False):
                txt = shape.text or ""
        except Exception:
            txt = ""
        # Tables: flatten rows into pipe-separated lines.
        try:
            if getattr(shape, "has_table", False):
                lines = []
                tbl = shape.table
                for row in tbl.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append((cell.text or "").strip())
                    lines.append(" | ".join([c for c in cells if c]))
                ttxt = "\n".join([l for l in lines if l.strip()])
                if ttxt:
                    txt = _join_nonempty(txt, ttxt, sep="\n")
        except Exception:
            pass
        return (txt or "").strip()

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        try:
            if slide.shapes.title is not None:
                title = slide.shapes.title.text or ""
        except Exception:
            title = ""

        body_parts: list[str] = []
        for shape in slide.shapes:
            t = _shape_text(shape)
            if t:
                body_parts.append(t)
        slide_text = "\n".join(body_parts).strip()

        notes_text = ""
        try:
            notes = slide.notes_slide
            # The notes slide often has multiple shapes; concatenate text frames.
            n_parts: list[str] = []
            for shape in notes.shapes:
                try:
                    if getattr(shape, "has_text_frame", False):
                        t = (shape.text or "").strip()
                        if t:
                            n_parts.append(t)
                except Exception:
                    continue
            notes_text = "\n".join(n_parts).strip()
        except Exception:
            notes_text = ""

        out.append(
            ExtractedSlide(
                slide_no=i,
                title=(title or "").strip(),
                slide_text=slide_text,
                notes_text=notes_text,
            )
        )

    return out


async def ingest_pptx_content_to_db(*, content_id: UUID) -> None:
    """
    Ingest a PPTX course content item into the Postgres retrieval layer.

    - Store per-slide combined text in document_pages (page_no = slide_no)
    - Store chunks in content_chunks with consistent page_start/page_end
    """
    settings = get_settings()
    if not getattr(settings, "rag_enabled", True):
        return
    if not settings.s3_bucket:
        return

    SessionLocal = get_session_maker()

    async with SessionLocal() as db:
        res = await db.execute(select(CourseContent).where(CourseContent.id == content_id))
        content = res.scalar_one_or_none()
        if content is None or not content.file_key:
            return
        if not _is_pptx(content):
            return

        course_id = content.course_id
        category = content.category
        file_key = content.file_key
        original_filename = content.original_filename

    s3 = _s3_client(settings)
    try:
        obj = s3.get_object(Bucket=settings.s3_bucket, Key=file_key)
        body = obj.get("Body")
        pptx_bytes = body.read() if body is not None else b""
    except Exception:
        return

    slides = _extract_pptx_slides(pptx_bytes)
    if not slides:
        return

    splitter = _text_splitter(settings)

    total_chars = 0
    for s in slides:
        total_chars += len((s.title or "")) + len((s.slide_text or "")) + len((s.notes_text or ""))
    low_text = total_chars < 200

    page_rows: list[DocumentPage] = []
    chunk_rows: list[ContentChunk] = []
    chunk_index = 0

    for s in slides:
        # Combine slide content with context and notes.
        header = f"Slide {s.slide_no}"
        if s.title:
            header = f"{header} — {s.title}"

        combined = _join_nonempty(
            header,
            s.slide_text,
            (f"Notes:\n{s.notes_text}" if s.notes_text else ""),
            sep="\n",
        )
        combined = combined.strip()

        page_rows.append(
            DocumentPage(
                course_id=course_id,
                content_id=content_id,
                page_no=int(s.slide_no),
                text=combined,
                text_sha256=_sha256_hex(combined),
            )
        )

        if not combined:
            continue

        chunks = splitter.split_text(combined)
        for c in chunks:
            c = (c or "").strip()
            if not c:
                continue
            meta: dict[str, Any] = {
                "doc_type": "slides",
                "source_kind": "pptx",
                "page_start": int(s.slide_no),
                "page_end": int(s.slide_no),
                "slide_no": int(s.slide_no),
                "original_filename": original_filename,
            }
            if s.title:
                meta["title"] = s.title
            if low_text:
                meta["extraction_warning"] = "low_text_extracted"
            chunk_rows.append(
                ContentChunk(
                    course_id=course_id,
                    content_id=content_id,
                    category=category,
                    chunk_index=chunk_index,
                    text=c,
                    meta=meta,
                )
            )
            chunk_index += 1

    async with SessionLocal() as db:
        await db.execute(delete(DocumentPage).where(DocumentPage.content_id == content_id))
        await db.execute(delete(ContentChunk).where(ContentChunk.content_id == content_id))
        if page_rows:
            db.add_all(page_rows)
        if chunk_rows:
            db.add_all(chunk_rows)
        await db.commit()


